from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)

def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=24, bold=False, alignment=PP_ALIGN.LEFT, color=BLACK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Consolas'
    p.alignment = alignment
    return tf

def add_bullet_text(text, tf, size=20, color=BLACK):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = 'Consolas'
    p.level = 0
    p.space_after = Pt(6)
    return p

def add_line(slide, x1, y1, x2, y2, color=BLACK):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x1), Inches(y1),
        Inches(x2 - x1), Inches(0.01)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ============================================================
# SLIDE 1 - Welcome
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1)

add_textbox(slide1, 1.5, 2.0, 10, 1.5, "PASSWORD AUDIT SYSTEM", size=44, bold=True, alignment=PP_ALIGN.CENTER)
add_line(slide1, 3.5, 3.5, 9.8, 3.5)
add_textbox(slide1, 1.5, 3.8, 10, 1, "Python Object-Oriented Programming Project", size=24, alignment=PP_ALIGN.CENTER)
add_textbox(slide1, 1.5, 4.8, 10, 0.8, "All 18 Syllabus Topics Demonstrated", size=20, alignment=PP_ALIGN.CENTER)
add_textbox(slide1, 1.5, 5.8, 10, 0.8, "41 Unit Tests - All Passing", size=20, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 - Goal
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2)

add_textbox(slide2, 0.8, 0.5, 12, 1, "GOAL", size=36, bold=True)
add_line(slide2, 0.8, 1.3, 12, 1.3)

tf = add_textbox(slide2, 0.8, 1.8, 11, 5, "", size=22)
tf.paragraphs[0].text = ""
goals = [
    "Load a file with multiple passwords (one per line)",
    "Run validation checks on each password",
    "  - Length check (min 8 characters)",
    "  - Complexity check (upper, lower, digit, special)",
    "  - Pattern check (against common weak passwords)",
    "Detect duplicate passwords",
    "Generate a full audit report file",
    "Demonstrate all 18 Python OOP syllabus topics",
]
for g in goals:
    add_bullet_text(g, tf, size=22)

# ============================================================
# SLIDE 3 - Tech Used
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3)

add_textbox(slide3, 0.8, 0.5, 12, 1, "TECH USED", size=36, bold=True)
add_line(slide3, 0.8, 1.3, 12, 1.3)

tf = add_textbox(slide3, 0.8, 1.8, 11, 5, "", size=22)
techs = [
    "Language: Python 3",
    "Libraries: abc (abstract base classes)",
    "Data Structures: List, Tuple, Dict, Set, Frozenset",
    "File I/O: read passwords, write report",
    "Exception Handling: try / except",
    "Testing: unittest module (41 tests)",
    "Documentation: DOCX generation (python-docx)",
]
for t in techs:
    add_bullet_text(t, tf, size=22)

# ============================================================
# SLIDE 4 - Features
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4)

add_textbox(slide4, 0.8, 0.5, 12, 1, "FEATURES", size=36, bold=True)
add_line(slide4, 0.8, 1.3, 12, 1.3)

# Left column
tf1 = add_textbox(slide4, 0.8, 1.8, 5.5, 5, "", size=20)
feats_left = [
    "6 Classes (OOP Architecture)",
    "Abstract base class (ABC)",
    "3 Check subclasses (Inheritance)",
    "Polymorphic run_check()",
    "Name-mangled encapsulation",
    "Conditional strength rating",
]
for f in feats_left:
    add_bullet_text(f, tf1, size=20)

# Right column
tf2 = add_textbox(slide4, 7, 1.8, 5.5, 5, "", size=20)
feats_right = [
    "Set-based duplicate detection",
    "Frozenset weak password list",
    "Dict {password: result} mapping",
    "List of tuples for report output",
    "Comprehension-based filtering",
    "Exception handling for files",
]
for f in feats_right:
    add_bullet_text(f, tf2, size=20)

add_line(slide4, 6.5, 1.8, 6.5, 6.5)

# ============================================================
# SLIDE 5 - Class Architecture
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5)

add_textbox(slide5, 0.8, 0.5, 12, 1, "CLASS ARCHITECTURE", size=36, bold=True)
add_line(slide5, 0.8, 1.3, 12, 1.3)

tf = add_textbox(slide5, 0.8, 1.8, 11, 5, "", size=20)
classes = [
    "PasswordCheck (ABC) - Abstract base",
    "  +-- LengthCheck      - Checks length >= 8",
    "  +-- ComplexityCheck  - Upper, lower, digit, special",
    "  +-- PatternCheck     - Frozenset weak passwords",
    "",
    "Password - Stores password + private __score, __issues",
    "PasswordAuditor - Loads file, detects dupes, orchestrates checks",
    "AuditReport - Generates formatted report file",
    "",
    "Topics: Inheritance, Polymorphism, Encapsulation, Abstraction",
]
for c in classes:
    add_bullet_text(c, tf, size=22 if c and not c.startswith("  ") else 20)

# ============================================================
# SLIDE 6 - End
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6)

add_textbox(slide6, 1.5, 2.5, 10, 1.5, "THANK YOU", size=44, bold=True, alignment=PP_ALIGN.CENTER)
add_line(slide6, 4.5, 3.8, 8.8, 3.8)
add_textbox(slide6, 1.5, 4.2, 10, 1, "Password Audit System", size=24, alignment=PP_ALIGN.CENTER)
add_textbox(slide6, 1.5, 5.0, 10, 0.8, "python password_audit.py", size=20, alignment=PP_ALIGN.CENTER)
add_textbox(slide6, 1.5, 5.6, 10, 0.8, "41/41 Tests Passing", size=18, alignment=PP_ALIGN.CENTER)

output_path = "Password_Audit_System.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
