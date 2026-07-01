import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_background(slide):
    # Apply solid dark charcoal background (RGB 20, 20, 20)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 20)

def add_slide_header(slide, title_text, category_text="CINEBOOK PRO"):
    # Add a thin gold category tag at the top
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.73), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = 'Arial'
    p_cat.font.size = Pt(9)
    p_cat.font.bold = True
    p_cat.font.color.rgb = RGBColor(228, 179, 99)  # Warm Amber Gold
    
    # Add the main slide title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.73), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255) # Pure White

def main():
    prs = Presentation()
    # Set to widescreen 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # =========================================================================
    # SLIDE 1: COVER
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1)
    
    # Title & Subtitle block
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0
    
    p1 = tf1.paragraphs[0]
    p1.text = "CINEBOOK PRO"
    p1.font.name = 'Arial'
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(228, 179, 99) # Gold Accent
    p1.space_after = Pt(12)
    
    p2 = tf1.add_paragraph()
    p2.text = "A Hybrid Multi-Language Movie Ticketing & Security Architecture"
    p2.font.name = 'Arial'
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.space_after = Pt(24)
    
    p3 = tf1.add_paragraph()
    p3.text = "Internship Project Presentation  |  Department of Cyber Security"
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(160, 160, 160)
    
    # Student Info block at bottom right
    info_box = slide1.shapes.add_textbox(Inches(7.5), Inches(5.2), Inches(4.8), Inches(1.5))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    tf_info.margin_left = tf_info.margin_right = tf_info.margin_top = tf_info.margin_bottom = 0
    p_info = tf_info.paragraphs[0]
    p_info.text = "Submitted by:"
    p_info.font.bold = True
    p_info.font.size = Pt(12)
    p_info.font.color.rgb = RGBColor(228, 179, 99)
    
    p_name = tf_info.add_paragraph()
    p_name.text = "Kolluri Chaturvedhi Narsimha\nHTNo: 24EG109A28\nAnurag University"
    p_name.font.size = Pt(12)
    p_name.font.color.rgb = RGBColor(220, 220, 220)

    # =========================================================================
    # SLIDE 2: TECH STACK & SYSTEM OVERVIEW
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2)
    add_slide_header(slide2, "Project Overview & Technology Stack")
    
    # Left column: Description
    desc_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = tf_desc.margin_right = tf_desc.margin_top = tf_desc.margin_bottom = 0
    
    p_overview = tf_desc.paragraphs[0]
    p_overview.text = "SYSTEM CONCEPT"
    p_overview.font.bold = True
    p_overview.font.size = Pt(14)
    p_overview.font.color.rgb = RGBColor(228, 179, 99)
    p_overview.space_after = Pt(8)
    
    p_body = tf_desc.add_paragraph()
    p_body.text = "CineBook Pro is designed to demonstrate multi-framework web orchestration coupled with native backend validation workflows. It mounts independent frontend frameworks side-by-side while utilizing robust multi-language backends to validate core ticketing details."
    p_body.font.size = Pt(13)
    p_body.font.color.rgb = RGBColor(220, 220, 220)
    p_body.space_after = Pt(20)
    
    p_core = tf_desc.add_paragraph()
    p_core.text = "KEY ATTRIBUTES"
    p_core.font.bold = True
    p_core.font.size = Pt(14)
    p_core.font.color.rgb = RGBColor(228, 179, 99)
    p_core.space_after = Pt(8)
    
    bullets = [
        "React-Vue frontend interoperability.",
        "Java-based security & memory visualization.",
        "Python password strength classification.",
        "AI-assisted prompt engineering development."
    ]
    for b in bullets:
        pb = tf_desc.add_paragraph()
        pb.text = "• " + b
        pb.font.size = Pt(12)
        pb.font.color.rgb = RGBColor(180, 180, 180)
        pb.space_after = Pt(6)

    # Right column: Tech Stack Grid
    grid_box = slide2.shapes.add_textbox(Inches(6.6), Inches(1.8), Inches(5.9), Inches(4.8))
    tf_grid = grid_box.text_frame
    tf_grid.word_wrap = True
    tf_grid.margin_left = tf_grid.margin_right = tf_grid.margin_top = tf_grid.margin_bottom = 0
    
    p_tech = tf_grid.paragraphs[0]
    p_tech.text = "TECHNOLOGY DETAILS"
    p_tech.font.bold = True
    p_tech.font.size = Pt(14)
    p_tech.font.color.rgb = RGBColor(228, 179, 99)
    p_tech.space_after = Pt(14)
    
    techs = [
        ("React (Host SPA)", "Manages global state, routing, and interactive ticket stubs."),
        ("Vue.js (Widgets)", "Encapsulates modular UI cards and password audit screens."),
        ("Java (Security Engine)", "Executes Luhn validations, logs file operations, and runs thread locks."),
        ("Python (Auditor)", "Runs structural validations and generates statistical plots."),
        ("Prompt Engineering", "Utilized AI prompts to debug code and optimize database locking.")
    ]
    for title, desc in techs:
        pt = tf_grid.add_paragraph()
        pt.text = title + "  —  "
        pt.font.bold = True
        pt.font.size = Pt(12)
        pt.font.color.rgb = RGBColor(255, 255, 255)
        
        run = pt.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = RGBColor(180, 180, 180)
        pt.space_after = Pt(12)

    # =========================================================================
    # SLIDE 3: SYSTEM ARCHITECTURE
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3)
    add_slide_header(slide3, "System Component Architecture")
    
    # Left: Image
    if os.path.exists("captures/system_architecture.png"):
        slide3.shapes.add_picture("captures/system_architecture.png", Inches(0.8), Inches(1.8), width=Inches(5.6))
    
    # Right: Bullet points
    details_box = slide3.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    tf_details = details_box.text_frame
    tf_details.word_wrap = True
    tf_details.margin_left = tf_details.margin_right = tf_details.margin_top = tf_details.margin_bottom = 0
    
    p_details_title = tf_details.paragraphs[0]
    p_details_title.text = "COMPONENT ARCHITECTURE"
    p_details_title.font.bold = True
    p_details_title.font.size = Pt(14)
    p_details_title.font.color.rgb = RGBColor(228, 179, 99)
    p_details_title.space_after = Pt(12)
    
    arch_points = [
        ("React Host & Vue Containers", "The React host utilizes hash routing and context managers while dynamically mounting Vue checkout widgets in the container DOM."),
        ("Java Security Middleware", "Implements Private encasing, card verification through Luhn's algorithm, transaction serialization via BufferedWriter, and concurrency checks."),
        ("Python Auditor & Terminal Console", "Derived password checking strategies from abstract parents, displaying evaluation stats on a dynamic Matplotlib chart."),
        ("File I/O & Storage Layer", "Tracks local session variables inside LocalStorage and stores credit details securely to disk in cards.dat.")
    ]
    for header, desc in arch_points:
        pt = tf_details.add_paragraph()
        pt.text = "• " + header + "\n"
        pt.font.bold = True
        pt.font.size = Pt(12)
        pt.font.color.rgb = RGBColor(255, 255, 255)
        pt.space_before = Pt(4)
        
        run = pt.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = RGBColor(170, 170, 170)
        pt.space_after = Pt(10)

    # =========================================================================
    # SLIDE 4: TRANSACTION SEQUENCE FLOW
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4)
    add_slide_header(slide4, "End-to-End Booking Sequence Flow")
    
    # Left: Image
    if os.path.exists("captures/booking_sequence_flow.png"):
        slide4.shapes.add_picture("captures/booking_sequence_flow.png", Inches(0.8), Inches(1.8), width=Inches(5.6))
        
    # Right: Phases list
    seq_box = slide4.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    tf_seq = seq_box.text_frame
    tf_seq.word_wrap = True
    tf_seq.margin_left = tf_seq.margin_right = tf_seq.margin_top = tf_seq.margin_bottom = 0
    
    p_seq_title = tf_seq.paragraphs[0]
    p_seq_title.text = "TRANSACTIONAL PHASES"
    p_seq_title.font.bold = True
    p_seq_title.font.size = Pt(14)
    p_seq_title.font.color.rgb = RGBColor(228, 179, 99)
    p_seq_title.space_after = Pt(12)
    
    phases = [
        ("Phase 1: Seating Lock", "A synchronized block launches Thread-0 on showtime selection, locking the coordinates for 5 minutes to avoid conflicts."),
        ("Phase 2: Card Verification", "The Luhn validation checksum checks bank card details, recording verified structures to disk via BufferedWriter."),
        ("Phase 3: Security Audit", "Validates password strengths. Java tracks simple length restrictions while Python computes entropy and displays charts."),
        ("Phase 4: Ticket Issue", "Generates interactive 3D ticket stubs with embedded SVG QR codes, triggering validation confirmations on click.")
    ]
    for header, desc in phases:
        pt = tf_seq.add_paragraph()
        pt.text = "• " + header + "\n"
        pt.font.bold = True
        pt.font.size = Pt(12)
        pt.font.color.rgb = RGBColor(255, 255, 255)
        pt.space_before = Pt(4)
        
        run = pt.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = RGBColor(170, 170, 170)
        pt.space_after = Pt(10)

    # =========================================================================
    # SLIDE 5: CODE INTEGRATION (React + Vue)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5)
    add_slide_header(slide5, "Technical Focus: React & Vue Interoperability")
    
    # Left: Explanation
    info_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    tf_info.margin_left = tf_info.margin_right = tf_info.margin_top = tf_info.margin_bottom = 0
    
    p_header = tf_info.paragraphs[0]
    p_header.text = "INTEROPERABILITY DETAILS"
    p_header.font.bold = True
    p_header.font.size = Pt(14)
    p_header.font.color.rgb = RGBColor(228, 179, 99)
    p_header.space_after = Pt(12)
    
    details = [
        ("React DOM Wrapper", "React coordinates the global session state and mounts a wrapper div container. It manages component visibility and transitions."),
        ("Vue Widget Mounts", "Inside React's lifecycle hooks, custom Vue apps are initialized and mounted to the ref wrappers. On unmounting, state cleanup is performed."),
        ("State Synchronization", "Data is bridged using custom browser event listeners, synchronizing updates in seat counts and invoices between React and Vue components.")
    ]
    for header, desc in details:
        pt = tf_info.add_paragraph()
        pt.text = "• " + header + "\n"
        pt.font.bold = True
        pt.font.size = Pt(12)
        pt.font.color.rgb = RGBColor(255, 255, 255)
        pt.space_before = Pt(6)
        
        run = pt.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = RGBColor(170, 170, 170)
        pt.space_after = Pt(10)
        
    # Right: Code block
    code_box = slide5.shapes.add_textbox(Inches(6.5), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_code = code_box.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = tf_code.margin_right = tf_code.margin_top = tf_code.margin_bottom = 0
    
    p_code_header = tf_code.paragraphs[0]
    p_code_header.text = "REACT LIFECYCLE MOUNTING CODE"
    p_code_header.font.bold = True
    p_code_header.font.size = Pt(14)
    p_code_header.font.color.rgb = RGBColor(228, 179, 99)
    p_code_header.space_after = Pt(14)
    
    code_p = tf_code.add_paragraph()
    code_p.text = (
        "// React component mounting Vue app\n"
        "const VueWrapper = () => {\n"
        "  const containerRef = useRef(null);\n\n"
        "  useEffect(() => {\n"
        "    // Create and mount Vue widget instance\n"
        "    const app = createApp(CardCheckout);\n"
        "    app.mount(containerRef.current);\n\n"
        "    return () => {\n"
        "      app.unmount(); // Unmount on clean up\n"
        "    };\n"
        "  }, []);\n\n"
        "  return <div ref={containerRef} />;\n"
        "};"
    )
    code_p.font.name = 'Courier New'
    code_p.font.size = Pt(11)
    code_p.font.color.rgb = RGBColor(240, 240, 240)

    # =========================================================================
    # SLIDE 6: JAVA SECURITY MODULES
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_background(slide6)
    add_slide_header(slide6, "Technical Focus: Java Security Layer")
    
    # Left: Code block
    code_box_java = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_code_java = code_box_java.text_frame
    tf_code_java.word_wrap = True
    tf_code_java.margin_left = tf_code_java.margin_right = tf_code_java.margin_top = tf_code_java.margin_bottom = 0
    
    p_code_java_header = tf_code_java.paragraphs[0]
    p_code_java_header.text = "LUHN ALGORITHM CARD CHECK"
    p_code_java_header.font.bold = True
    p_code_java_header.font.size = Pt(14)
    p_code_java_header.font.color.rgb = RGBColor(228, 179, 99)
    p_code_java_header.space_after = Pt(14)
    
    code_java_p = tf_code_java.add_paragraph()
    code_java_p.text = (
        "public boolean checkLuhn(String cardNo) {\n"
        "    int nDigits = cardNo.length();\n"
        "    int nSum = 0; boolean isSecond = false;\n"
        "    for (int i = nDigits - 1; i >= 0; i--) {\n"
        "        int d = cardNo.charAt(i) - '0';\n"
        "        if (isSecond == true) d = d * 2;\n"
        "        nSum += d / 10;\n"
        "        nSum += d % 10;\n"
        "        isSecond = !isSecond;\n"
        "    }\n"
        "    return (nSum % 10 == 0);\n"
        "}"
    )
    code_java_p.font.name = 'Courier New'
    code_java_p.font.size = Pt(11)
    code_java_p.font.color.rgb = RGBColor(240, 240, 240)
    
    # Right: Java Details
    java_details_box = slide6.shapes.add_textbox(Inches(7.3), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_java_details = java_details_box.text_frame
    tf_java_details.word_wrap = True
    tf_java_details.margin_left = tf_java_details.margin_right = tf_java_details.margin_top = tf_java_details.margin_bottom = 0
    
    p_java_details_title = tf_java_details.paragraphs[0]
    p_java_details_title.text = "CORE JAVA IMPL"
    p_java_details_title.font.bold = True
    p_java_details_title.font.size = Pt(14)
    p_java_details_title.font.color.rgb = RGBColor(228, 179, 99)
    p_java_details_title.space_after = Pt(12)
    
    java_points = [
        ("Luhn Check Algorithm", "Executes mod-10 double-every-second digit sum validations on inputs to prevent fake/invalid card numbers."),
        ("OOP Encapsulation Pattern", "Strictly keeps attributes private, accessing them via clean getter/setter objects to block arbitrary edits."),
        ("ArrayList Allocation", "Stores validated instances to memory dynamically, mapping object heap locations to logs."),
        ("File System Logging", "Streams logs out to cards.dat on disk using Java's BufferedWriter class, securing records.")
    ]
    for header, desc in java_points:
        pt = tf_java_details.add_paragraph()
        pt.text = "• " + header + "\n"
        pt.font.bold = True
        pt.font.size = Pt(12)
        pt.font.color.rgb = RGBColor(255, 255, 255)
        pt.space_before = Pt(4)
        
        run = pt.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = RGBColor(170, 170, 170)
        pt.space_after = Pt(10)

    # =========================================================================
    # SLIDE 7: PYTHON AUDIT MODULES
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_background(slide7)
    add_slide_header(slide7, "Technical Focus: Python Audit Layer")
    
    # Left: Explanation
    py_info_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_py_info = py_info_box.text_frame
    tf_py_info.word_wrap = True
    tf_py_info.margin_left = tf_py_info.margin_right = tf_py_info.margin_top = tf_py_info.margin_bottom = 0
    
    p_py_header = tf_py_info.paragraphs[0]
    p_py_header.text = "AUDIT SPECIFICATION"
    p_py_header.font.bold = True
    p_py_header.font.size = Pt(14)
    p_py_header.font.color.rgb = RGBColor(228, 179, 99)
    p_py_header.space_after = Pt(12)
    
    py_details = [
        ("Polymorphic Abstract Structure", "Utilizes abstract classes defining baseline check functions. Concrete sub-classes overwrite checking parameters."),
        ("Entropy & Pattern Scanning", "Applies regular expressions to scan passwords for repeated sequences, length rules, and character sets."),
        ("Matplotlib Graph Generation", "Saves audit results and vulnerability scores as a bar chart, saving a PNG which is rendered in the app's HTML canvas.")
    ]
    for header, desc in py_details:
        pt = tf_py_info.add_paragraph()
        pt.text = "• " + header + "\n"
        pt.font.bold = True
        pt.font.size = Pt(12)
        pt.font.color.rgb = RGBColor(255, 255, 255)
        pt.space_before = Pt(6)
        
        run = pt.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = RGBColor(170, 170, 170)
        pt.space_after = Pt(10)
        
    # Right: Python code block
    py_code_box = slide7.shapes.add_textbox(Inches(6.5), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_py_code = py_code_box.text_frame
    tf_py_code.word_wrap = True
    tf_py_code.margin_left = tf_py_code.margin_right = tf_py_code.margin_top = tf_py_code.margin_bottom = 0
    
    p_py_code_header = tf_py_code.paragraphs[0]
    p_py_code_header.text = "ABSTRACT AUDITING ENGINE"
    p_py_code_header.font.bold = True
    p_py_code_header.font.size = Pt(14)
    p_py_code_header.font.color.rgb = RGBColor(228, 179, 99)
    p_py_code_header.space_after = Pt(14)
    
    py_code_p = tf_py_code.add_paragraph()
    py_code_p.text = (
        "from abc import ABC, abstractmethod\n\n"
        "class PasswordCheck(ABC):\n"
        "    @abstractmethod\n"
        "    def validate(self, pwd: str) -> bool:\n"
        "        pass\n\n"
        "class LengthCheck(PasswordCheck):\n"
        "    def validate(self, pwd: str) -> bool:\n"
        "        return len(pwd) >= 8\n\n"
        "class ComplexityCheck(PasswordCheck):\n"
        "    def validate(self, pwd: str) -> bool:\n"
        "        # RegEx checks for lower/upper/digit/special\n"
        "        return bool(re.search(r'[A-Z]', pwd))"
    )
    py_code_p.font.name = 'Courier New'
    py_code_p.font.size = Pt(11)
    py_code_p.font.color.rgb = RGBColor(240, 240, 240)

    # =========================================================================
    # SLIDE 8: INTERACTIVE SCREENSHOTS
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_background(slide8)
    add_slide_header(slide8, "Application Features & Interface Designs")
    
    # 4-image Grid Layout
    # Top-Left: Home
    if os.path.exists("captures/screenshots/01_home_screen.png"):
        slide8.shapes.add_picture("captures/screenshots/01_home_screen.png", Inches(0.8), Inches(1.8), width=Inches(2.7))
    # Top-Right: Seat Map
    if os.path.exists("captures/screenshots/05_seat_map_selected_recommended.png"):
        slide8.shapes.add_picture("captures/screenshots/05_seat_map_selected_recommended.png", Inches(3.7), Inches(1.8), width=Inches(2.7))
    # Bottom-Left: Ticket
    if os.path.exists("captures/screenshots/13_ticket_3d_modal_unscanned.png"):
        slide8.shapes.add_picture("captures/screenshots/13_ticket_3d_modal_unscanned.png", Inches(0.8), Inches(4.3), width=Inches(2.7))
    # Bottom-Right: Python Audit
    if os.path.exists("captures/screenshots/16_python_audit_results.png"):
        slide8.shapes.add_picture("captures/screenshots/16_python_audit_results.png", Inches(3.7), Inches(4.3), width=Inches(2.7))
        
    # Text Details
    grid_details_box = slide8.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    tf_grid_details = grid_details_box.text_frame
    tf_grid_details.word_wrap = True
    tf_grid_details.margin_left = tf_grid_details.margin_right = tf_grid_details.margin_top = tf_grid_details.margin_bottom = 0
    
    p_grid_details_title = tf_grid_details.paragraphs[0]
    p_grid_details_title.text = "CORE USER CAPABILITIES"
    p_grid_details_title.font.bold = True
    p_grid_details_title.font.size = Pt(14)
    p_grid_details_title.font.color.rgb = RGBColor(228, 179, 99)
    p_grid_details_title.space_after = Pt(12)
    
    grid_details = [
        ("Home & Showtime Selector", "Frosted glass visual layouts dynamically loading show details. Responsive grids align items smoothly."),
        ("3D Curved Seating Map", "Interactive seats mapped to custom hooks. Visualizes selected seats, recommended paths, and price categories."),
        ("Parallax 3D Ticket Modal", "Features responsive card tilts. Users can hover to examine validation details and click to verify QR passes."),
        ("Password Auditing Dashboard", "Simulates server security terminal overlays. Computes entropy levels and compiles charts of weak parameters.")
    ]
    for header, desc in grid_details:
        pt = tf_grid_details.add_paragraph()
        pt.text = "• " + header + "\n"
        pt.font.bold = True
        pt.font.size = Pt(12)
        pt.font.color.rgb = RGBColor(255, 255, 255)
        pt.space_before = Pt(4)
        
        run = pt.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = RGBColor(170, 170, 170)
        pt.space_after = Pt(10)

    # =========================================================================
    # SLIDE 9: CHALLENGES & RESOLUTIONS
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_background(slide9)
    add_slide_header(slide9, "Technical Challenges & Engineering Mitigations")
    
    challenges_box = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.8))
    tf_challenges = challenges_box.text_frame
    tf_challenges.word_wrap = True
    tf_challenges.margin_left = tf_challenges.margin_right = tf_challenges.margin_top = tf_challenges.margin_bottom = 0
    
    p_chal_title = tf_challenges.paragraphs[0]
    p_chal_title.text = "ENGINEERING MITIGATIONS"
    p_chal_title.font.bold = True
    p_chal_title.font.size = Pt(14)
    p_chal_title.font.color.rgb = RGBColor(228, 179, 99)
    p_chal_title.space_after = Pt(14)
    
    chals = [
        ("Double-Booking Race Condition", 
         "Challenge: Simultaneous booking attempts by multiple user contexts created double-booking conflicts.\n"
         "Mitigation: Implemented Mutex seating locks. Thread-0 acquires validation limits via Java synchronized blocks, preventing multi-user overlaps."),
        ("React-Vue Interoperability", 
         "Challenge: Mounting independent reactive frameworks side-by-side inside DOM trees produced component lifecycle collisions.\n"
         "Mitigation: Managed DOM wrappers with React useRef and loaded Vue apps inside useEffect hooks. Cleaned memory caches on component unmounting."),
        ("Validation Exception Propagation", 
         "Challenge: Formatting errors in deep class libraries failed to notify frontend layers.\n"
         "Mitigation: Structured custom exceptions that trace errors to sideboards, displaying error logs directly in client consoles.")
    ]
    for title, desc in chals:
        pt = tf_challenges.add_paragraph()
        pt.text = "• " + title + "\n"
        pt.font.bold = True
        pt.font.size = Pt(12.5)
        pt.font.color.rgb = RGBColor(255, 255, 255)
        
        run = pt.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(11.5)
        run.font.color.rgb = RGBColor(170, 170, 170)
        pt.space_after = Pt(14)

    # =========================================================================
    # SLIDE 10: CONCLUSION & FUTURE DIRECTIONS
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_background(slide10)
    add_slide_header(slide10, "Conclusion & Learnings")
    
    # Left: Key learnings
    concl_box = slide10.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_concl = concl_box.text_frame
    tf_concl.word_wrap = True
    tf_concl.margin_left = tf_concl.margin_right = tf_concl.margin_top = tf_concl.margin_bottom = 0
    
    p_concl_title = tf_concl.paragraphs[0]
    p_concl_title.text = "CORE INTERNSHIP LEARNINGS"
    p_concl_title.font.bold = True
    p_concl_title.font.size = Pt(14)
    p_concl_title.font.color.rgb = RGBColor(228, 179, 99)
    p_concl_title.space_after = Pt(14)
    
    learnings = [
        ("Software Engineering Design", "Learned how to split responsibilities across independent, specialized layers (Web UI, Security logic, and Auditor modules)."),
        ("Cross-Framework Integration", "Gained deep understanding of DOM lifecycles by mounting, bridging, and cleaning memory parameters between React and Vue."),
        ("Robust OOP & Concurrency", "Gained practical experience using Java synchronized locks, encapsulation patterns, custom exception structures, and list serialization."),
        ("AI-Assisted Prompting Flow", "Learned to leverage prompt engineering techniques to draft mock parameters, speed up debugging, and format logs.")
    ]
    for title, desc in learnings:
        pt = tf_concl.add_paragraph()
        pt.text = "• " + title + "\n"
        pt.font.bold = True
        pt.font.size = Pt(12)
        pt.font.color.rgb = RGBColor(255, 255, 255)
        pt.space_before = Pt(4)
        
        run = pt.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = RGBColor(170, 170, 170)
        pt.space_after = Pt(8)

    # Right: Thank you / Q&A block
    qa_box = slide10.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8))
    tf_qa = qa_box.text_frame
    tf_qa.word_wrap = True
    tf_qa.margin_left = tf_qa.margin_right = tf_qa.margin_top = tf_qa.margin_bottom = 0
    
    p_qa_title = tf_qa.paragraphs[0]
    p_qa_title.text = "QUESTIONS & ANSWERS"
    p_qa_title.font.bold = True
    p_qa_title.font.size = Pt(14)
    p_qa_title.font.color.rgb = RGBColor(228, 179, 99)
    p_qa_title.space_after = Pt(24)
    
    p_thank = tf_qa.add_paragraph()
    p_thank.text = "Thank You."
    p_thank.font.bold = True
    p_thank.font.size = Pt(36)
    p_thank.font.color.rgb = RGBColor(255, 255, 255)
    p_thank.space_after = Pt(12)
    
    p_contact = tf_qa.add_paragraph()
    p_contact.text = "Feel free to ask any technical questions regarding the hybrid architecture, card check validations, password auditing algorithms, or prompt engineering methods."
    p_contact.font.size = Pt(13)
    p_contact.font.color.rgb = RGBColor(200, 200, 200)

    # Save presentation
    output_filename = "CineBook_Pro_Internship_Presentation.pptx"
    prs.save(output_filename)
    print(f"SUCCESS: Presentation saved as '{output_filename}'")

if __name__ == '__main__':
    main()
